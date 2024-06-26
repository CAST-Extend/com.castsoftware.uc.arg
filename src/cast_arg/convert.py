from cast_arg.restCall import AipData,HLData
#from cast_arg.powerpoint import PowerPoint
from cast_arg.actionPlan import ActionPlan
from cast_arg.config import Config
from cast_arg.pages.hl_cloud import CloudMaturity
from cast_arg.pages.hl_greenIt import GreenIt
from cast_arg.pages.hl_summary import HighlightSummary
from cast_arg.pages.hl_benchmark import HighlightBenchmark
from cast_arg.pages.mri_strengh_improvement import StrengthImprovment
from cast_arg.pages.mri_grades import MRIGrades
from cast_arg.pages.mri_sizing import MRISizing
from cast_arg.pages.mri_tech_detail_table import TechDetailTable

from cast_arg.stats import OssStats,AIPStats,LicenseStats



from cast_common.mri import MRI
from cast_common.logger import Logger,DEBUG, INFO, WARN
from cast_common.util import find_nth, no_dups, list_to_text,format_table
from cast_arg.powerpoint import PowerPoint
from cast_common.highlight import Highlight

from copy import deepcopy

from pandas import DataFrame
from pptx import Presentation
from pptx.dml.color import RGBColor
from IPython.display import display
from os import getcwd
from os.path import abspath,dirname,exists
from site import getsitepackages

import pandas as pd
import numpy as np 
import math
import json
import datetime

__author__ = "Nevin Kaplan"
__email__ = "n.kaplan@castsoftware.com"
__copyright__ = "Copyright 2022, CAST Software"

class GeneratePPT(Logger):
    _ppt = None
    _aip_data = None
    _effort_df = None
    _output_folder = None
    _project_name = None
    _template = None

    _generate_HL = _hl_base_url = _hl_user = hl_pswd = None
    _generate_AIP = _aip_base_url = _aip_user = _aip_pswd = None
    _hl_instance = None
    _hl_apps_df = pd.DataFrame()
    _hl_app_list = []

    day_rate=1600

    def __init__(self, config:Config):
        super().__init__("generate",config.logging_generate)
        self._config = config

        out = abspath(f"{config.output}/Project {config.project} - Tech DD Findings.pptx")
        self.out=out
        self.info(f'Generating {out}')

        self._ppt = PowerPoint(config)

        # TODO: Handle cases where on HL data is needed and not AIP.
        self.hl_pages = []
        if config.aip_active:
            self.info("Collecting AIP Data")
            try:
                self._aip_data = AipData(config,log_level=config.logging_aip)
            except PermissionError:
                self.error('Invalid MRI REST API Credentials!')
                exit (1)

            self.mri_pages = [
                MRIGrades(log_level=INFO,ppt=self._ppt),
                MRISizing(),
                TechDetailTable(),
                StrengthImprovment()
            ]
        if config.hl_active:
            self.info("Collecting Highlight Data")
            hl = Highlight(hl_base_url=config.hl_url,hl_user=config.hl_user,hl_pswd=config.hl_password, \
                           hl_instance=config.hl_instance,hl_apps=config.hl_list)
            self.hl_portfolio_pages = [
                HighlightSummary(self.day_rate,self._config.output,ppt=self._ppt)
            ]
            self.hl_pages = [
                CloudMaturity(),
                GreenIt(),
                HighlightSummary(self.day_rate),
                HighlightBenchmark()
            ]
            self._hl_data = HLData(config,log_level=config.logging_highlight)

        #project level work
        app_cnt = len(config.application)

        # self.remove_proc_slides(self._generate_procs)

        if app_cnt == 1: 
            s = self._ppt.get_shape_by_name('port_level_slide')
            pass
            #self._ppt.delete_slide
        else:
            self._ppt.duplicate_slides(app_cnt)
            self._ppt.copy_block("each_app",["app"],app_cnt)
            # self._ppt.save()
            # return 

        self._ppt.replace_text("{app_per_page}","",tbd_for_blanks=False)

        self.expand_tables(config,['project_overview'])
        self.replace_all_text()

    def expand_tables(self,config:Config,table_names:list):
        app_cnt = len(config.application)
        for tbl_name in table_names:
            try:
                table = self._ppt.get_shape_by_name(tbl_name)
                if table is None:
                    raise ValueError(f'Table not found in template: {tbl_name}')
                max_rows = self._ppt.table_max_rows(table) 
                if max_rows < app_cnt:
                    max_rows = app_cnt
                table = src_table = table.table

                for idx in range(1,app_cnt):
                    new_row = deepcopy(src_table._tbl.tr_lst[1]) 

                    #spill over to the next page
                    if idx+1 == max_rows:
                        table = self._ppt.get_shape_by_name(f'{tbl_name}_spill')
                        if table is None:
                            raise ValueError(f'Table not found in template: {tbl_name}_spill')
                            if table is None:
                                raise ValueError(f'Table not found in template: {tbl_name}_spill')
                        table = table.table

                    table._tbl.append(new_row)
                    c_row = len(table.rows)-1

                    for col in range(0,len(table.columns)):
                        new_cell = table.cell(c_row,col) 
                        for p in new_cell.text_frame.paragraphs:
                            self._ppt._replace_paragraph_text(p, '{app1',f'{{app{idx+1}')
                        pass
                    pass
            except ValueError:
                self.warning(f'table not found {tbl_name}')

    def remove_proc_slides(self,keep_it):
        indexes=[]
        for idx, slide in enumerate(self._ppt._prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text == "{proc_analysis}":
                            self._ppt.replace_slide_text(slide,"{proc_analysis}","")
                            indexes.append(idx)

        for idx in reversed(indexes):
            if not keep_it:
                self._ppt.delete_slide(idx)

    # def save_ppt(self):
    #     while True:
    #         try:
    #             self._ppt.save()
    #             self.info(f'{self.out} saved.')
    #             return 
    #         except PermissionError: 
    #             while answer := input (f'Error writing {self.out} powerpoint document, Retry [Y or N]:'):
    #                 if answer.upper() == 'Y':
    #                     break
    #                 elif answer.upper() == 'N':
    #                     return
    #                 else:
    #                     continue

    def replace_all_text(self):
        app_cnt = len(self._config.application)
        self._ppt.replace_text("{project}", self._config.project)
        self._ppt.replace_text("{app_count}",app_cnt)
        self._ppt.replace_text("{company}",self._config.company)

        mydate = datetime.datetime.now()
        month = mydate.strftime("%B")
        year = f'{mydate.year} '
        self._ppt.replace_text("{month}",month)
        self._ppt.replace_text("{year}",year)

        # replace AIP data global to all applications
        if self._config.aip_active:
            all_apps_avg_grade = self._aip_data.calc_grades_all_apps()
#            self._ppt.replace_text("{all_apps}",self._aip_data.get_all_app_text())
            self._ppt.replace_risk_factor(all_apps_avg_grade,search_str="{summary_")
            risk_grades = self._aip_data.calc_health_grades_high_risk(all_apps_avg_grade)
            if risk_grades.empty:
                risk_grades = self._aip_data.calc_health_grades_medium_risk(all_apps_avg_grade)
            self._ppt.replace_text("{summary_at_risk_factors}",list_to_text(risk_grades.index.tolist()).lower())
        
            if risk_grades.empty:
                self._ppt.replace_block('{immediate_action}','{end_immediate_action}','')
            else:
                self._ppt.replace_text("{immediate_action}","") 
                self._ppt.replace_text("{end_immediate_action}","") 
                self._ppt.replace_text("{high_risk_grade_names}",self._aip_data.text_from_list(risk_grades.index.values.tolist()))

        app_list = self._config.aip_list
        hl_list = self._config.hl_list

        day_rate = self.day_rate

        summary_near_term=AIPStats(day_rate,logger_level=self._config.logging_generate)
        summary_fix_now=AIPStats(day_rate,logger_level=self._config.logging_generate)
        summary_mid_long_term=AIPStats(day_rate,logger_level=self._config.logging_generate)
        hl_summary=OssStats('',day_rate,logger_level=self._config.logging_generate)
        hl_summary_critical=OssStats('',day_rate,logger_level=self._config.logging_generate)
        hl_summary_high_near=OssStats('',day_rate,logger_level=self._config.logging_generate)
        lic_summary = LicenseStats(logger_level=self._config.logging_generate)
        summary_components = 0

        """Highlight Portfolio level pages """
        if self._config.hl_active:
            for proc in self.hl_portfolio_pages:
                proc.report(hl_list)

        self._ppt.replace_text('{app_cnt}',app_cnt)

        from inflect import engine
        self._ppt.replace_text('{app_cnt_as_word}',engine().number_to_words(app_cnt))

        for idx in range(0,app_cnt):
            # create instance of action plan class 
            self.ap = ActionPlan (app_list,self._aip_data,self._ppt,self._config.output,day_rate)
            fix_now_total=AIPStats(day_rate,logger_level=self._config.logging_generate)
            near_term_total=AIPStats(day_rate,logger_level=self._config.logging_generate)
            mid_long_term=AIPStats(day_rate,logger_level=self._config.logging_generate)
            summary_total=AIPStats(day_rate,logger_level=self._config.logging_generate)

            hl_near_term_total=AIPStats(day_rate,logger_level=self._config.logging_generate)

            # replace application specific AIP data
            app_no = idx+1
            app_id = app_list[idx]
            hl_id = hl_list[idx]
            app_title = self._config.title_list[idx]
            self.info(f'********************* Working on pages for {app_title} ******************************')
            self._ppt.replace_text(f'{{app{app_no}_name}}',app_title)


            if self._config.aip_active:
                if self._aip_data.has_data(app_id):
                    self.info('Preparing AIP Data')

                    #Run MRI pages
                    for proc in self.mri_pages:
                        self.info(f'Generating {proc.description}')
                        proc.run(app_id,app_no)

                    self.info('Filling risk factors for the executive summary page')
                    # do risk factors for the executive summary page
                    risk_grades = self.each_risk_factor(self._aip_data,app_id, app_no)
                    self._ppt.replace_text(f'{{app{app_no}_high_risk_grade_names}}',list_to_text(risk_grades.index.values))

                    snapshot = self._aip_data.snapshot(app=app_id)
                    self._ppt.replace_text(f'{{app{app_no}_all_technogies}}',list_to_text(snapshot['technology']))
                    
                    """
                        Populate the document insites page
                        The necessary data is found in the loc_tbl

                        This section fetches data for Documentation slide, excludes certain columns, sorts with few column, 
                        and based on score of particular element seggregates colors(Red, Yellow & Green) and update to app1_doc_table element.
                        255,168,168 - Red shades
                        255,234,168 - Yellow shades
                        168,228,195 - Green shades
                    """
                    doc_df = self._aip_data.doc_compliance(app_id).copy()
                    doc_df.drop(columns=['Key','Total','Weight'],inplace=True) #'Detail',
                    doc_df.sort_values(by=['Score','Rule'], inplace=True)
                    doc_df['RGB'] = np.where(doc_df.Score >= 3,'194,236,213',np.where(doc_df.Score < 2,'255,210,210','255,240,194'))
                    doc_df.Score = doc_df.Score.map('{:.2f}'.format)
                    self._ppt.update_table(f'app{app_no}_doc_table',doc_df,app_id,include_index=False,background='RGB')
                    self.fill_critical_rules(app_id,app_no)

                    """
                        Get Action Plan data and combine it for the various slide combinations

                        Executive Summary uses a combination of extreme and high violation data
                        Action Plan midigation slice is divided into three sections
                            1. Fix before sigining contains only extreme violation data
                            2. Near term contains high violation data
                            3. Long term contains medium and low violation data
                            4. Excecutive Summary contains fix now and high violation data
                    """
                    self.ap.fill_action_plan(app_id,app_no)
                    fix_now_total.add_effort(self.ap.fix_now.effort)
                    fix_now_total.add_violations(self.ap.fix_now.violations)

                    near_term_total.add_effort(self.ap.high.effort)
                    near_term_total.add_violations(self.ap.high.violations)

                    mid_long_term.add_effort(self.ap.medium.effort)
                    mid_long_term.add_violations(self.ap.medium.violations)
                    mid_long_term.add_data(self.ap.medium.data)
                    mid_long_term.add_effort(self.ap.low.effort)
                    mid_long_term.add_violations(self.ap.low.violations)
                    mid_long_term.add_data(self.ap.low.data)

                    # summary_near_term.add_effort(self.ap.high.effort)
                    # summary_near_term.add_effort(self.ap.medium.effort)
                    # summary_near_term.add_violations(self.ap.high.violations)
                    # summary_near_term.add_violations(self.ap.medium.violations)

                    """
                        ISO-5055 slide
                            - most of the work is being done in the restCall class
                            - use iso_rules to retrieve the data
                            - add it to the app1_iso5055 table in the template
                    """
                    iso_df = self._aip_data.iso_rules(app_id)
                    if not iso_df.empty:
                        try:
                            iso_df.loc[iso_df['violation']=='','background']='205,218,226'
                            iso_df.loc[iso_df['violation']!='','background']='255,255,255'
                            self._ppt.update_table(f'app{app_no}_iso5055',iso_df,app_id,
                                                include_index=False,background='background')
                                        
                            pourcentage_iso5055 = iso_df["count"].sum()
                            
                            iso_Maintainaility = iso_df[iso_df.category == 'Maintainability' ]
                            iso_MaintainailityCall = iso_Maintainaility["count"].sum()
                            self._ppt.replace_text(f'{{app{app_no}_ISO_5055}}', round((iso_MaintainailityCall/(pourcentage_iso5055/2))*100,1))
                        except ValueError as ex:
                            self.warning(ex)                    
            #replaceHighlight application specific data
            if self._config.hl_active and self._hl_data.has_data(hl_id):
                for proc in self.hl_pages:
                    proc.report(hl_id,app_no)

                try:
                    (oss_crit,oss_high,oss_med,lic,components) = self.oss_risk_assessment(hl_id,app_no,day_rate)
                    fix_now_total.add_effort(oss_crit.effort)
                    fix_now_total.add_violations(oss_crit.violations)

                    near_term_total.add_effort(oss_high.effort)
                    near_term_total.add_effort(oss_med.effort)

                    hl_near_term_total.add_effort(oss_high.effort)
                    hl_near_term_total.add_effort(oss_med.effort)
                    hl_near_term_total.add_violations(oss_high.violations)
                    hl_near_term_total.add_violations(oss_med.violations)
                    hl_near_term_total.replace_text(self._ppt,app_no,'hl_near_term_total')

                    hl_summary_critical.add_components(oss_crit.components)
                    hl_summary_critical.add_violations(oss_crit.violations)

                    hl_summary_high_near.add_components(oss_high.components)
                    hl_summary_high_near.add_components(oss_med.components)

                    hl_summary.add_components(components)

                    lic_summary.add_high(lic.high)
                    lic_summary.add_medium(lic.medium)
                    lic_summary.add_low(lic.low)


                except KeyError as ex:
                    self.warning(f'OSS information not found {str(ex)}')


                """
                    Cloud ready excel sheet generation
                """
                # try:
                #     cloud = self._hl_data.get_cloud_info(hl_id)
                #     cloud = cloud[['cloudRequirement.display','Technology','cloudRequirement.ruleType','cloudRequirement.criticality','contributionScore','roadblocks']]
                #     file_name = f'{self._config.output}/cloud-{self._config.title_list[idx]}.xlsx'
                #     writer = pd.ExcelWriter(file_name, engine='xlsxwriter')
                #     col_widths=[50,10,10,10,10,10,10]
                #     cloud_tab = format_table(writer,cloud,'Cloud Data',col_widths)
                #     writer.close()
                # except Exception as e:
                #     self.error(f'unknown error while processing cloud ready data: {str(e)}')


            summary_fix_now.add_effort(fix_now_total.effort)
            summary_fix_now.add_violations(fix_now_total.violations)
            summary_near_term.add_effort(near_term_total.effort)
            summary_near_term.add_violations(near_term_total.violations)
            summary_mid_long_term.add_effort(mid_long_term.effort)
            summary_mid_long_term.add_violations(mid_long_term.violations)

            summary_total.add_effort(summary_fix_now.effort)
            summary_total.add_effort(summary_near_term.effort)
            summary_total.add_effort(summary_mid_long_term.effort)
            summary_total.add_violations(summary_fix_now.violations)
            summary_total.add_violations(summary_near_term.violations)
            summary_total.add_violations(summary_mid_long_term.violations)
            summary_total.add_data(summary_fix_now.data)
            summary_total.add_data(summary_near_term.data)
            summary_total.add_data(summary_mid_long_term.data)
            summary_total.replace_text(self._ppt,app_no,'summary_total')

            #replace text for all aip and HL statistics in the powerpoint document
            self.ap.fix_now.replace_text(self._ppt,app_no,'fix_now')
            self.ap.high.replace_text(self._ppt,app_no,'high')
            self.ap.medium.replace_text(self._ppt,app_no,'medium')
            self.ap.low.replace_text(self._ppt,app_no,'low')

            mid_long_term.replace_text(self._ppt,app_no,'mid_long_term')

            fix_now_total.replace_text(self._ppt,app_no,'fix_now_total')
            near_term_total.replace_text(self._ppt,app_no,'near_term_total')

            # short_term.replace_text(self._ppt,app_no,'st')
            # long_term.replace_text(self._ppt,app_no,'lt')
        
        #summary_fix_now_aip_hl.replace_text(self._ppt,'','hl_summary_fix_now')
        summary_fix_now.replace_text(self._ppt,'','summary_fix_now')
        summary_near_term.replace_text(self._ppt,'','summary_near_term')

        hl_summary_critical.replace_text(self._ppt,'_summary_crit')
        hl_summary_high_near.replace_text(self._ppt,'_summary_high_near')
        lic_summary.replace_text(self._ppt,'_summary')
        hl_summary.replace_text(self._ppt,'_summary')
        self._ppt.replace_text('{app_summary_comp_tot}',summary_components)

        show_stopper_flg = 'some'
        if self.ap.fix_now.effort == 0:
            show_stopper_flg = 'no'
        self._ppt.replace_text('{show_stopper_flg}',show_stopper_flg)  

        avg_cost = float(summary_fix_now.cost/app_cnt)   
        if avg_cost < 50:
            hml_cost_flg = 'low'
        elif avg_cost > 50 and avg_cost < 100:
            hml_cost_flg = 'medium'
        else:
            hml_cost_flg = 'high'
        self._ppt.replace_text('{hml_cost_flag}',hml_cost_flg)  

        self._ppt.replace_text('{daily_rate}',self.ap._day_rate)  

    def each_risk_factor(self, aip_data, app_id, app_no):
        # collect the high risk grades
        # if there are no high ris grades then use medium risk
        app_level_grades = aip_data.get_app_grades(app_id)
        risk_grades = aip_data.calc_health_grades_high_risk(app_level_grades)
        risk_catagory="high"
        if risk_grades.empty:
            risk_catagory="medium"
            risk_grades = aip_data.calc_health_grades_medium_risk(app_level_grades)
        
        #in the event all health risk factors are low risk
        if risk_grades.empty:
            self._ppt.replace_block(f'{{app{app_no}_risk_detail}}',
                            f'{{end_app{app_no}_risk_detail}}',
                            "no high-risk health factors")
        else: 
            rpl_str = f'{{app{app_no}_risk_category}}'
            self._ppt.replace_text(rpl_str,risk_catagory)
            self.debug(f'replaced {rpl_str} with {risk_catagory}')

            self._ppt.copy_block(f'app{app_no}_each_risk_factor',["_risk_name","_risk_grade"],len(risk_grades.count(axis=1)))
            f=1
            for index, row in risk_grades.T.items():
                rpl_str = f'{{app{app_no}_risk_name{f}}}'
                self._ppt.replace_text(rpl_str,index)
                self.debug(f'replaced {rpl_str} with {index}')

                rpl_str = f'{{app{app_no}_risk_grade{f}}}'
                value = row['All'].round(2)
                self._ppt.replace_text(rpl_str,value)
                self.debug(f'replaced {rpl_str} with {value}')
                f=f+1

        self._ppt.remove_empty_placeholders()
        return risk_grades

    def fill_critical_rules(self,app_id,app_no):
        self.info('Filling critical rules table')
        rules_df = self._aip_data.critical_rules(app_id)
        if not rules_df.empty:
            rules_df = rules_df[['rulePattern.name','rulePattern.critical']]
            rule_summary_df=rules_df.groupby(['rulePattern.name']).size().reset_index(name='counts').sort_values(by=['counts'],ascending=False)
            rule_summary_df=rule_summary_df.head(5)
            self._ppt.update_table(f'app{app_no}_top_violations',rule_summary_df,app_id,include_index=False)
        else:
            self.warning('This application contains no critical violations')

    def oss_risk_assessment(self,hl_id,app_no,day_rate):
        self.info('Filling OSS risk assessment table')
        lic_df=self._hl_data.get_lic_info(hl_id)
        lic_df=self._hl_data.sort_lic_info(lic_df)
        oss_df=self._hl_data.get_cve_info(hl_id)
        lic = LicenseStats(logger_level=self._config.logging_generate)
        # lic_summary = pd.DataFrame(columns=['License Type','Risk Factor','Component Count','Example'])

        oss_crit = OssStats(hl_id,day_rate,self._hl_data,'crit',logger_level=self._config.logging_generate)
        oss_crit_comp_tot = self._hl_data.get_cve_crit_comp_tot(hl_id)

        oss_high = OssStats(hl_id,day_rate,self._hl_data,'high',logger_level=self._config.logging_generate)
        oss_high_comp_tot = self._hl_data.get_cve_high_comp_tot(hl_id)

        oss_med = OssStats(hl_id,day_rate,self._hl_data,'med',logger_level=self._config.logging_generate)
        oss_med_comp_tot = self._hl_data.get_cve_med_comp_tot(hl_id)

        total_components = self._hl_data.get_oss_cmpn_tot(hl_id)

        oss_crit.replace_text(self._ppt,app_no)
        oss_high.replace_text(self._ppt,app_no)
        oss_med.replace_text(self._ppt,app_no)

        self._ppt.replace_text(f'{{app{app_no}_oss_cmpn_tot}}',total_components)
        oss_df = oss_df[(oss_df['critical']!='') | (oss_df['high']!='')]
        if not oss_df.empty:
            tbl_name = f'app{app_no}_hl_table_cve'
            try:
                self._ppt.update_table(tbl_name,oss_df,hl_id,include_index=False)
            except ValueError as ex:
                self.warning(ex)
        self.info('Filling OSS license table')
        '''
            License compliance table
        
        '''
        lic_summary=self._hl_data.get_lic_info(hl_id)
        if not lic_summary.empty:
            lic_summary=lic_summary[['component','version','release','license','risk']].drop_duplicates()
            lic_summary.sort_values(['risk','license','component','version'],inplace=True)

            lic_summary = lic_summary.groupby(['risk','license'])['component'].apply(lambda x: ','.join(x)).reset_index()
            lic_summary['comp count']=lic_summary['component'].str.count(',')+1
            
            #remove duplicates but show full count
            lic_summary['component']=lic_summary['component'].map(lambda x: no_dups(x,',',True))
            
            #only show the first 5 components
            lic_summary['component']=lic_summary['component'].map(lambda x: x[:find_nth(x,',',6)])
            lic_summary=lic_summary[['license','risk','comp count','component']]
            lic_summary.sort_values(['risk','license','comp count'],inplace=True)

            #remove low and undefined records from the table
            lic_summary=lic_summary[lic_summary["risk"].str.contains("Low")==False]
            lic_summary=lic_summary[lic_summary["risk"].str.contains("Undefined")==False]

            # are there any records left?
            if not lic_summary.empty:
                #modify the forground color
                lic_summary.loc[lic_summary['risk']=='High','forground']='211,76,76'
                lic_summary.loc[lic_summary['risk']=='Medium','forground']='127,127,127'

                #update the powerpoint table
                tbl_name = f'app{app_no}_hl_table_lic_risks'
                try:
                    self._ppt.update_table(tbl_name,lic_summary,hl_id, include_index=False)
                except ValueError as ex:
                    self.warning(ex)
            
                lic.high = lic_summary[lic_summary['risk']=='High']['comp count'].sum()
                lic.medium = lic_summary[lic_summary['risk']=='Medium']['comp count'].sum()
                lic.low = lic_summary[lic_summary['risk']=='Low']['comp count'].sum()

                # #add the high and medium license risk counts to the deck
                # self._ppt.replace_text(f'{{app{app_no}_high_lic_tot}}',
                #     lic_summary[lic_summary['risk']=='High']['comp count'].sum())
                # self._ppt.replace_text(f'{{app{app_no}_med_lic_tot}}',
                #     lic_summary[lic_summary['risk']=='Medium']['comp count'].sum())

        lic.replace_text(self._ppt,app_no)
        return (oss_crit,oss_high,oss_med,lic,total_components)





