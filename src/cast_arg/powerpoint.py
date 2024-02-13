from cast_common.powerpoint import PowerPoint as common_ppt
from cast_common.logger import INFO
from cast_arg.config import Config
from pptx.chart.data import CategoryChartData
from pandas import Series,DataFrame
from os.path import abspath

__author__ = "Nevin Kaplan"
__email__ = "n.kaplan@castsoftware.com"
__copyright__ = "Copyright 2022, CAST Software"

def yes_no_response(msg:str) -> bool:
    while answer := input (f'{msg},  [Y or N]?'):
        if answer.upper() == 'Y':
            return True
        elif answer.upper() == 'N':
            return False
        else:
            continue

def numeric_response(msg:str,value:int,max_rows) -> int:
    msg_str=msg
    while True:
        first = True
        answer = input(f'{msg_str}: ')
        if len(answer)==0:
            answer = f'{value}'
        if answer.isnumeric():
            answer = int(answer)
            if answer > max_rows:
                msg_str = f'{msg} [Too many rows, data contains {max_rows} rows]'
                continue
            return answer
        else:
                msg_str = f'{msg} [Input must be numberic]'
                first = False
            
        

class PowerPoint(common_ppt):
    """
    This module holds techniques explicitly designed for 
    constructing PowerPoint slides. The majority of its capabilities 
    come from the PowerPoint module in the com.castsoftware.uc.python.common library. 

    Args:
        common_ppt: PowerPoint class found in com.castsoftware.uc.python.common
    """

#    def __init__(self,template:str=None,output:str=None,log_level=INFO) -> None:
    def __init__(self,config:Config,log_level=INFO) -> None:
        super().__init__(config.template,config.output,log_level)
        self._out = abspath(f"{config.output}/Project {config.project} - Tech DD Findings.pptx")        
        self.config=config

    def save(self):
        while True:
            try:
                self._prs.save(self._out)
                self.log.info(f'{self._out} saved.')
                return 
            except PermissionError as pe: 
                if not yes_no_response(f'Error writing {self._out} powerpoint document, Retry'):
                    return 
            except Exception as ex:
                self.log.error(f'General Exception while saving PowerPoint document: {ex}')
                raise ex


    def update_table(self,name,data:DataFrame, app:str, include_index=True,interactive:bool=False):
        self.log.info(f'Updating table {name}')

        """
            is there a tables entry in the config.json? if not add it now
            check if the table name is in the config.json tables array
                if not add it now
            finally set a local variable to table
        """
        if name not in self.config.tables:
            self.config.tables[name]={}
        table = self.config.tables[name]
            
        table_shape = self.get_shape_by_name(name)
        if table_shape is None:
            raise ValueError(f'Table not found in template: {name}')

        last_row = self.table_max_rows(table_shape)
        if 'row' not in table:
            table['row']=last_row
            self.config.save()
        elif last_row != table['row']:
            last_row = table['row']

        #work on main table
        while True:
            super().update_table(name,data,max_rows=last_row,include_index=include_index)
            if interactive:
                self.log.info(f'Deck {name} table for application {app} has been updated.')
                if rsp := numeric_response(f'The current row number of lines is {last_row}, enter a new value or RETURN to accept',last_row,len(data)):
                    if rsp == last_row:
                        break
                    last_row = rsp
                if table['row']!=last_row:
                    table['row']=last_row
                    self.config.save()
            else:
                break
            
        self.replace_text(f'{{name}}',last_row)

        #now do spillover table
        try:
            spill_data = data.iloc[last_row:]
            spill_name = f'{name}_spill_1'
            table_shape = self.get_shape_by_name(spill_name)
            if table_shape is None:
                raise ValueError(f'Table not found in template: {spill_name}')

            super().update_table(spill_name,spill_data,max_rows=len(spill_data),include_index=include_index)
        except ValueError as ve:
            self.log.warning(ve)

        self.save()
        pass

    def replace_risk_factor(self, grades:Series, app_no:int=0, search_str:str=None):
        if search_str == None:
            search_str=f'{{app{app_no}_risk_'
        
        for slide in self._prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.find(search_str)!=-1:
                            run = self._merge_runs(paragraph)
                            # cur_text=''
                            # first=True
                            # for run in paragraph.runs:
                            #     cur_text = cur_text + run.text
                            #     if first != True:
                            #         self.delete_run(run)
                            #     first=False
                            # run = paragraph.runs[0]
                            # run.text = cur_text

                            while True:
                                t = cur_text.find(search_str)
                                if t == -1: 
                                    break
                                g = cur_text[t+len(search_str):]
                                g = g[:g.find("}")]
                                if g:
                                    try:
                                        grade=grades[g]
                                        risk = ''
                                        if grade < 2:
                                            risk = 'high'
                                        elif grade < 3:
                                            risk = 'medium'
                                        else:
                                            risk = 'low'

                                        cur_text = cur_text.replace(f'{search_str}{g}}}',risk)
                                        run.text = cur_text
                                    except KeyError:
                                        self.debug(f'invalid key: {g}')
                                        break;

    def update_grade_slider(self, shape,data):
        #shape = self.get_shape_by_name(name)
        try:
            if shape.has_chart:
                chart_data = CategoryChartData()
                # chart_data.categories = titles
                
                chart_data.categories = ['grade']
                chart_data.add_series('Series 1', data)
                shape.chart.replace_data(chart_data)
        except AttributeError as ae:
            self.warning(f'Attribute Error, invalid template configuration: {shape.name} ({ae})')
        except KeyError as ke:
            self.warning(f'Key Error, invalid template configuration: {shape.name} ({ke})')
        except TypeError as t:
            self.warning(f'Type Error, invalid template configuration: {shape.name} ({t})')

    def remove_empty_placeholders(self):
        for slide in self._prs.slides:
            for placeholder in slide.shapes.placeholders:
                if placeholder.has_text_frame and placeholder.text_frame.text == "":
                    sp = placeholder._sp
                    sp.getparent().remove(sp)

    def replace_loc(self, loc, app_no):
        loc_short = "{0:,.0f} LoC".format(loc) 
        if loc > 1000000:
            loc_short = "~{0:,.2f} MLoC".format(loc/1000000) 
        elif loc < 1000000 and loc > 1000:
            loc_short = "~{0:,.0f} KLoC".format(loc/1000) 
        self.replace_text(f'{{app{app_no}_loc}}',f'{loc:,.0f}')
        self.replace_text(f'{{app{app_no}_loc_short}}',loc_short)

        size_catagory = 'small'
        if loc > 1000000:
            size_catagory = 'very large'
        elif loc > 500000:
            size_catagory = 'large'
        self.replace_text(f'{{app{app_no}_loc_category}}',size_catagory)

