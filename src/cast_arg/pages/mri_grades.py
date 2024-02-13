from cast_arg.pages.mri_report import MRIPage
from cast_arg.powerpoint import PowerPoint
from cast_common.util import format_table,list_to_text

from pandas import Series
from pptx.dml.color import RGBColor

class MRIGrades(MRIPage):
    description = 'Calculating MRI Grades'

    def report(self,app_name:str,app_no:int) -> bool:

        app_level_grades = self.get_app_grades(app_name)
        for name, value in app_level_grades.T.items():
            # fill grades
            grade = round(value,2)
            rpl_str = f'{{app{app_no}_grade_{name}}}'
            self.ppt.replace_text(rpl_str,grade)
            self._log.debug(f'replaced {rpl_str} with {grade}')

            # fill grade risk factor (high, medium or low)
            rpl_str = f'{{app{app_no}_risk_{name}}}'
            risk = ''
            if grade < 2:
                risk = 'high'
            elif grade < 3:
                risk = 'medium'
            else:
                risk = 'low'
            self.ppt.replace_text(rpl_str,risk)
            self._log.debug(f'replaced {rpl_str} with {risk}')

            #update grade box color and slider postion 
            id_base = f'app{app_no}_grade'
            box_name = f'{id_base}_{name}_box'
            txt_name = f'{id_base}_{name}_text'
            slider_name = f'{id_base}_{name}_slider'
            color = self.get_grade_color(grade)

            for slide in self.ppt._prs.slides:
                box = self.ppt.get_shape_by_name(box_name,slide)
                if not box is None:
                    box.line.color.rgb = color

                txt = self.ppt.get_shape_by_name(txt_name,slide)
                if not txt is None and txt.has_text_frame:
                    paragraphs = txt.text_frame.paragraphs
                    self.ppt.change_paragraph_color(paragraphs[0],color)

                slider = self.ppt.get_shape_by_name(slider_name,slide)
                if not slider is None:
                    self.ppt.update_grade_slider(slider,[grade])
        
        #calculate high and medium risk factors
        risk_grades = self.calc_health_grades_high_risk(app_level_grades)
        if risk_grades.empty:
            risk_grades = self.calc_health_grades_medium_risk(app_level_grades)
        self.ppt.replace_text(f'{{app{app_no}_at_risk_grade_names}}',list_to_text(risk_grades.index.tolist()).lower())
        self.replace_risk_factors(app_level_grades,app_no)

    def replace_risk_factors(self,grades:Series,app_no:int):

        for key in grades.keys():
            grade=grades[key]
            if grade < 2:
                risk = 'high'
            elif grade < 3:
                risk = 'medium'
            else:
                risk = 'low'

            self.ppt.replace_text(f'{{app{app_no}_risk_{key}}}',risk)
            pass

        pass 

    def get_grade_color(self,grade):
        rgb = 0
        if grade > 3:
            rgb = RGBColor(0,176,80) # light green
        elif grade <3 and grade > 2:
            rgb = RGBColor(214,142,48) # yellow
        else:
            rgb = RGBColor(255,0,0) # red
        return rgb
