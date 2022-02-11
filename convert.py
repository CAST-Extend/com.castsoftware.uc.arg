from logging import DEBUG, info, warn
from pandas.core.frame import DataFrame
from restCall import AipRestCall, AipData, HLRestCall, HLData
from pptx import Presentation
from powerpoint import PowerPoint
from jproperties import Properties 
from pptx.dml.color import RGBColor
from actionPlan import ActionPlan
from logger import Logger
from IPython.display import display
from config import Config
from util import find_nth, no_dups



import pandas as pd
import numpy as np 
import util 
import math
import argparse
import json
import datetime

class GeneratePPT(Logger):
    _ppt = None
    _aip_data = None
    _effort_df = None
    _output_folder = None
    _project_name = None
    _template = None

    _generate_HL = _hl_base_url = _hl_user = hl_pswd = None
    _generate_AIP = _aip_base_url = _aip_user = _aip_pswd = None
    _hl_instance = None
    _hl_apps_df = pd.DataFrame()
    _hl_app_list = []

    def __init__(self, config):
        super().__init__()
        out = f"{config.output}/Project {config.project} - Tech DD Findings.pptx"
        self.info(f'Generating {out}')


        # TODO: Handle cases where on HL data is needed and not AIP.
        if config.aip_active:
            self.info("Collecting AIP Data")
            self._aip_data = AipData(config)
        if config.hl_active:
            self.info("Collecting Highlight Data")
            self._hl_data = HLData(config)

        self._ppt = PowerPoint(config.template, out)
        #project level work
        app_cnt = len(config.application)

        # self.remove_proc_slides(self._generate_procs)

        self._ppt.duplicate_slides(app_cnt)
        self._ppt.copy_block("each_app",["app"],app_cnt)

        self._ppt.replace_text("{app#_","{app1_")

        self.replace_all_text(config)

    def remove_proc_slides(self,keep_it):
        indexes=[]
        for idx, slide in enumerate(self._ppt._prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text == "{proc_analysis}":
                            self._ppt.replace_slide_text(slide,"{proc_analysis}","")
                            indexes.append(idx)

        for idx in reversed(indexes):
            if not keep_it:
                self._ppt.delete_slide(idx)

    def save_ppt(self):
        self._ppt.save()

    def replace_all_text(self,config):
        app_cnt = len(config.application)
        self._ppt.replace_text("{project}", config.project)
        self._ppt.replace_text("{app_count}",app_cnt)
        self._ppt.replace_text("{company}",config.company)

        mydate = datetime.datetime.now()
        month = mydate.strftime("%B")
        year = f'{mydate.year} '
        self._ppt.replace_text("{month}",month)
        self._ppt.replace_text("{year}",year)

        # replace AIP data global to all applications
        if config.aip_active:
            all_apps_avg_grade = self._aip_data.calc_grades_all_apps()
#            self._ppt.replace_text("{all_apps}",self._aip_data.get_all_app_text())
            self._ppt.replace_risk_factor(all_apps_avg_grade,search_str="{summary_")
            risk_grades = self._aip_data.calc_health_grades_high_risk(all_apps_avg_grade)
            if risk_grades.empty:
                risk_grades = self._aip_data.calc_health_grades_medium_risk(all_apps_avg_grade)
            self._ppt.replace_text("{summary_at_risk_factors}",util.list_to_text(risk_grades.index.tolist()).lower())
        
            if risk_grades.empty:
                self._ppt.replace_block('{immediate_action}','{end_immediate_action}','')
            else:
                self._ppt.replace_text("{immediate_action}","") 
                self._ppt.replace_text("{end_immediate_action}","") 
                self._ppt.replace_text("{high_risk_grade_names}",self._aip_data.text_from_list(risk_grades.index.values.tolist()))

        app_list = config.aip_list
        hl_list = config.hl_list
        title_list = config.title_list

        # create instance of action plan class 
        ap = ActionPlan (app_list,self._aip_data,self._ppt,config.output)

        summary_total_cost = 0

        for idx in range(0,app_cnt):
            # replace application specific AIP data
            app_no = idx+1
            app_id = app_list[idx]
            hl_id = hl_list[idx]
            app_title = title_list[idx]
            self.info(f'Working on pages for {app_title}')
            self._ppt.replace_text(f'{{app{app_no}_name}}',app_title)

            aip_fix_now_eff = aip_fix_now_cost = aip_fix_now_vio_cnt = 0
            aip_near_term_eff = aip_near_term_cost = aip_near_term_vio_cnt = 0
            aip_mid_eff = aip_mid_cost = aip_mid_vio_cnt = 0
            aip_long_term_eff = aip_long_term_cost = aip_long_term_vio_cnt = 0

            aip_fix_now_eff = aip_fix_now_cost = aip_fix_now_vio_cnt = 0
            aip_near_term_eff = aip_near_term_cost = aip_near_term_vio_cnt = 0
            aip_mid_eff = aip_mid_cost = aip_mid_vio_cnt = 0
            aip_long_term_eff = aip_long_term_cost = aip_long_term_vio_cnt = 0

            fix_now_eff = fix_now_cost = fix_now_vio_cnt = 0

            fix_now_data = DataFrame()
            fix_now_bus_txt = ''
            fix_now_vio_txt = ''

            near_term_eff = near_term_cost = near_term_vio_cnt = 0
            near_term_data = DataFrame()
            near_term_bus_txt = ''
            near_term_vio_txt = ''

            mid_eff = mid_cost = mid_vio_cnt = 0
            mid_data = DataFrame()
            mid_bus_txt = ''
            mid_vio_txt = ''

            # low_eff = low_cost = low_vio_cnt = 0
            # low_data = DataFrame()
            # low_bus_txt = ''
            # low_vio_txt = ''

            long_term_eff = long_term_cost = long_term_vio_cnt = 0
            long_term_data = DataFrame()
            long_term_bus_txt = ''
            long_term_vio_txt = ''

            summary_eff = summary_cost = summary_vio_cnt = 0
            summary_data = DataFrame()
            summary_bus_txt = ''
            summary_vio_txt = ''

            if config.aip_active:
                if self._aip_data.has_data(app_id):
                    #self.info(f'Working on {app_id} ({self._appl_title[app_id]})')

                    # do risk factors for the executive summary page
                    risk_grades = util.each_risk_factor(self._ppt, self._aip_data,app_id, app_no)
                    self._ppt.replace_text(f'{{app{app_no}_high_risk_grade_names}}',util.list_to_text(risk_grades.index.values))


                    # Technical Overview - Technical details TABLE
                    grade_all = self._aip_data.get_app_grades(app_id)
                    self._ppt.replace_risk_factor(grade_all,app_no)
                    grade_by_tech_df = self._aip_data.get_grade_by_tech(app_id)
                    grades = grade_by_tech_df.drop(['Documentation',"ISO","ISO_EFF","ISO_MAINT","ISO_REL","ISO_SEC"],axis=1)
                    self._ppt.update_table(f'app{app_no}_grade_by_tech_table',grades)

                    # Technical Overview - Lines of code by technology GRAPH
                    self._ppt.update_chart(f'app{app_no}_sizing_pie_chart',grade_by_tech_df['LOC'])

                    snapshot = self._aip_data.snapshot(app=app_id)
                    # app_name = snapshot['name']
                    # self._ppt.replace_text(f'{{app{app_no}_name}}',app_name)
                    self._ppt.replace_text(f'{{app{app_no}_all_technogies}}',util.list_to_text(snapshot['technology']))

                    #calculate high and medium risk factors
                    risk_grades = self._aip_data.calc_health_grades_high_risk(grade_all)
                    if risk_grades.empty:
                        risk_grades = self._aip_data.calc_health_grades_medium_risk(grade_all)
                    self._ppt.replace_text(f'{{app{app_no}_at_risk_grade_names}}',util.list_to_text(risk_grades.index.tolist()).lower())


                    """
                        Populate the strengths and improvement page
                        The necessary data is found in the loc_tbl
                    """
                    imp_df = self._aip_data.tqi_compliance(app_id)
                    imp_df.drop(columns=['Weight','Total','Succeeded','Compliance'],inplace=True)
                    imp_df.sort_values(by=['Score','Rule'], inplace=True, ascending = False)

                    file_name = f'{config.output}/health-{title_list[idx]}.xlsx'
                    writer = pd.ExcelWriter(file_name, engine='xlsxwriter')
                    col_widths=[50,50,10,10,10]
                    cloud_tab = util.format_table(writer,imp_df,'Health Data',col_widths)
                    writer.save()

                    imp_df.drop(columns=['Detail'],inplace=True)
                    imp_df['RGB'] = np.where(imp_df.Score >= 3,'194,236,213',\
                        np.where(imp_df.Score < 2,'255,210,210','255,240,194'))
                    imp_df.Score = imp_df.Score.map('{:.2f}'.format)

                    imp_df['Cause']=''
                    with open('tech.json') as json_file:
                        tech_data = json.load(json_file)
                    imp_df['Cause']=imp_df['Key'].map(tech_data)

                    imp_df = imp_df[['Rule','Score','Cause','Failed','RGB']]
                    self._ppt.update_table(f'app{app_no}_imp_table',imp_df,include_index=False,background='RGB')

                    """
                        Populate the document insites page
                        The necessary data is found in the loc_tbl

                        This section fetches data for Documentation slide, excludes certain columns, sorts with few column, 
                        and based on score of particular element seggregates colors(Red, Yellow & Green) and update to app1_doc_table element.
                        255,168,168 - Red shades
                        255,234,168 - Yellow shades
                        168,228,195 - Green shades
                    """
                    doc_df = self._aip_data.doc_compliance(app_id)
                    doc_df.drop(columns=['Key','Total','Weight'],inplace=True) #'Detail',
                    doc_df.sort_values(by=['Score','Rule'], inplace=True)
                    doc_df['RGB'] = np.where(doc_df.Score >= 3,'194,236,213',np.where(doc_df.Score < 2,'255,210,210','255,240,194'))
                    doc_df.Score = doc_df.Score.map('{:.2f}'.format)
                    self._ppt.update_table(f'app{app_no}_doc_table',doc_df,include_index=False,background='RGB')
                    
                    loc_df = self._aip_data.get_loc_sizing(app_id)
                    loc = loc_df['Number of Code Lines']
                    self._ppt.replace_loc(loc,app_no)

                    loc_tbl = pd.DataFrame.from_dict(data=self._aip_data.get_loc_sizing(app_id),orient='index').drop('Critical Violations')
                    loc_tbl = loc_tbl.rename(columns={0:'loc'})
                    loc_tbl['percent'] = round((loc_tbl['loc'] / loc_tbl['loc'].sum()) * 100,2)
                    loc_tbl['loc']=pd.Series(["{0:,.0f}".format(val) for val in loc_tbl['loc']], index = loc_tbl.index)

                    percent_comment = loc_tbl.loc['Number of Comment Lines','percent']
                    percent_comment_out = loc_tbl.loc['Number of Commented-out Code Lines','percent']

                    if percent_comment < 15:
                        comment_level='low'
                    elif percent_comment > 15 and percent_comment <= 20:
                        comment_level='good'
                    else:
                        comment_level='high'
                    
                    self._ppt.replace_text(f'{{app{app_no}_comment_hl}}',comment_level)
                    self._ppt.replace_text(f'{{app{app_no}_comment_level}}',comment_level)
                    self._ppt.replace_text(f'{{app{app_no}_comment_pct}}',percent_comment)
                    self._ppt.replace_text(f'{{app{app_no}_comment_out_pct}}',percent_comment_out)

                    loc_tbl['percent']=pd.Series(["{0:.2f}%".format(val) for val in loc_tbl['percent']], index = loc_tbl.index)
                    self._ppt.update_table(f'app{app_no}_loc_table',loc_tbl,has_header=False)
                    self._ppt.update_chart(f'app{app_no}_loc_pie_chart',loc_tbl['loc'])

                    self._ppt.replace_grade(grade_all,app_no)

                    self.fill_sizing(app_id,app_no)
                    self.fill_critical_rules(app_id,app_no)


                    """
                        Get Action Plan data and combine it for the various slide combinations

                        Executive Summary uses a combination of extreme and high violation data
                        Action Plan midigation slice is divided into three sections
                            1. Fix before sigining contains only extreme violation data
                            2. Near term contains high violation data
                            3. Long term contains medium and low violation data
                            4. Excecutive Summary contains fix now and high violation data
                    """
                    ap.fill_action_plan(app_id,app_no)
                    self.fill_violations(app_id,app_no)

                    #calculate and print aip fix now action plan midigation items
                    (aip_fix_now_eff, aip_fix_now_cost, aip_fix_now_vio_cnt, aip_fix_now_data) = \
                        ap.get_extreme_costing()
                    aip_fix_now_bus_txt = util.list_to_text(ap.business_criteria(aip_fix_now_data)) + ' '
                    aip_fix_now_vio_txt = ap.list_violations(aip_fix_now_data)

                    self._ppt.replace_text(f'{{app{app_no}_aip_fn_eff}}',aip_fix_now_eff)
                    self._ppt.replace_text(f'{{app{app_no}_aip_fn_cost}}',aip_fix_now_cost)
                    self._ppt.replace_text(f'{{app{app_no}_aip_fn_vio_cnt}}',aip_fix_now_vio_cnt)
                    self._ppt.replace_text(f'{{app{app_no}_aip_fn_bus_txt}}',aip_fix_now_bus_txt)
                    self._ppt.replace_text(f'{{app{app_no}_aip_fn_vio_txt}}',aip_fix_now_vio_txt)

                    (aip_near_term_eff, aip_near_term_cost, aip_near_term_vio_cnt, aip_near_term_data) = \
                        ap.get_high_costing()
                    aip_near_term_bus_txt = util.list_to_text(ap.business_criteria(aip_near_term_data)) + ' '
                    aip_near_term_vio_txt = ap.list_violations(aip_near_term_data)
                    self._ppt.replace_text(f'{{app{app_no}_aip_nt_eff}}',aip_near_term_eff)
                    self._ppt.replace_text(f'{{app{app_no}_aip_nt_cost}}',aip_near_term_cost)
                    self._ppt.replace_text(f'{{app{app_no}_aip_nt_vio_cnt}}',aip_near_term_vio_cnt)
                    self._ppt.replace_text(f'{{app{app_no}_aip_nt_bus_txt}}',aip_near_term_bus_txt)
                    self._ppt.replace_text(f'{{app{app_no}_aip_nt_vio_txt}}',aip_near_term_vio_txt)


                    (aip_mid_term_eff, aip_mid_term_cost, aip_mid_term_vio_cnt, aip_mid_term_data) = ap.get_med_costing()
                    aip_mid_term_bus_txt = util.list_to_text(ap.business_criteria(aip_mid_term_data)) + ' '
                    aip_mid_term_vio_txt = ap.list_violations(aip_mid_term_data)
                    self._ppt.replace_text(f'{{app{app_no}_aip_mt_eff}}',aip_mid_term_eff)
                    self._ppt.replace_text(f'{{app{app_no}_aip_mt_cost}}',aip_mid_term_cost)
                    self._ppt.replace_text(f'{{app{app_no}_aip_mt_vio_cnt}}',aip_mid_term_vio_cnt)
                    self._ppt.replace_text(f'{{app{app_no}_aip_mt_bus_txt}}',aip_mid_term_bus_txt)
                    self._ppt.replace_text(f'{{app{app_no}_aip_mt_vio_txt}}',aip_mid_term_vio_txt)

                    (aip_low_eff, aip_low_cost, aip_low_vio_cnt, aip_low_data) = ap.get_low_costing()
                    aip_low_bus_txt = util.list_to_text(ap.business_criteria(aip_low_data)) + ' '
                    aip_low_vio_txt = ap.list_violations(aip_low_data)
                    self._ppt.replace_text(f'{{app{app_no}_aip_low_eff}}',aip_low_eff)
                    self._ppt.replace_text(f'{{app{app_no}_aip_low_cost}}',aip_low_cost)
                    self._ppt.replace_text(f'{{app{app_no}_aip_low_vio_cnt}}',aip_low_vio_cnt)
                    self._ppt.replace_text(f'{{app{app_no}_aip_low_bus_txt}}',aip_low_bus_txt)
                    self._ppt.replace_text(f'{{app{app_no}_aip_low_vio_txt}}',aip_low_vio_txt)

                    aip_long_term_data = pd.concat([aip_low_data,aip_mid_term_data],ignore_index=True)
                    aip_long_term_bus_txt = util.list_to_text(ap.business_criteria(aip_long_term_data)) + ' '
                    aip_long_term_vio_txt = ap.list_violations(aip_long_term_data)
                    aip_long_term_eff = int(mid_eff) + int(aip_low_eff)
                    aip_long_term_cost = float(mid_cost) + float(aip_low_cost)
                    aip_long_term_vio_cnt = int(mid_vio_cnt) + int(aip_low_vio_cnt)
                    self._ppt.replace_text(f'{{app{app_no}_aip_lt_eff}}',aip_long_term_eff)
                    self._ppt.replace_text(f'{{app{app_no}_aip_lt_cost}}',aip_long_term_cost)
                    self._ppt.replace_text(f'{{app{app_no}_aip_lt_vio_cnt}}',aip_long_term_vio_cnt)
                    self._ppt.replace_text(f'{{app{app_no}_aip_lt_bus_txt}}',aip_long_term_bus_txt)
                    self._ppt.replace_text(f'{{app{app_no}_aip_lt_vio_txt}}',aip_long_term_vio_txt)

                    summary_data = pd.concat([aip_fix_now_data,aip_near_term_data],ignore_index=True)
                    summary_bus_txt = util.list_to_text(ap.business_criteria(summary_data)) + ' '
                    summary_vio_txt = ap.list_violations(summary_data)
                    summary_eff = int(aip_fix_now_eff) + int(aip_near_term_eff)
                    summary_cost = float(aip_fix_now_cost) + float(aip_near_term_cost)
                    summary_vio_cnt = int(aip_fix_now_vio_cnt) + int(aip_near_term_vio_cnt)

                    fix_now_eff = int(aip_fix_now_eff) 
                    fix_now_cost = round(aip_fix_now_cost,2)
                    near_term_eff = int(aip_near_term_eff)
                    near_term_cost = round(aip_near_term_cost,2)

                    """
                        ISO-5055 slide
                            - most of the work is being done in the restCall class
                            - use iso_rules to retrieve the data
                            - add it to the app1_iso5055 table in the template
                    """
                    iso_df = self._aip_data.iso_rules(app_id)
                    iso_df.loc[iso_df['violation']=='','background']='205,218,226'
                    iso_df.loc[iso_df['violation']!='','background']='255,255,255'
                    self._ppt.update_table(f'app{app_no}_iso5055',iso_df,
                                           include_index=False,background='background')
            #replaceHighlight application specific data
            if config.hl_active and self._hl_data.has_data(hl_id):
                lic_df=self._hl_data.get_lic_info(hl_id)
                lic_df=self._hl_data.sort_lic_info(lic_df)
                oss_df=self._hl_data.get_cve_info(hl_id)
                # lic_summary = pd.DataFrame(columns=['License Type','Risk Factor','Component Count','Example'])

                crit_cve = self._hl_data.get_cve_crit_tot(hl_id)
                crit_comp_tot = self._hl_data.get_cve_crit_comp_tot(hl_id)

                high_cve = self._hl_data.get_cve_high_tot(hl_id)
                high_comp_tot = self._hl_data.get_cve_high_comp_tot(hl_id)

                med_cve = self._hl_data.get_cve_med_tot(hl_id)
                med_comp_tot = self._hl_data.get_cve_med_comp_tot(hl_id)

                oss_cmpnt_tot = self._hl_data.get_oss_cmpn_tot(hl_id)

                self._ppt.replace_text(f'{{app{app_no}_crit_sec_tot}}',crit_cve)
                self._ppt.replace_text(f'{{app{app_no}_high_sec_tot}}',high_cve)
                self._ppt.replace_text(f'{{app{app_no}_med_sec_tot}}',med_cve)

                self._ppt.replace_text(f'{{app{app_no}_crit_cve_comp_ct}}',crit_comp_tot)
                self._ppt.replace_text(f'{{app{app_no}_high_cve_comp_ct}}',high_comp_tot)
                self._ppt.replace_text(f'{{app{app_no}_med_cve_comp_ct}}',med_comp_tot)

                self._ppt.replace_text(f'{{app{app_no}_oss_cmpn_tot}}',oss_cmpnt_tot)

                if crit_cve is None:
                    crit_cve_eff = 0
                    crit_cve_cost = 0
                else:
                    crit_cve_eff = math.ceil(crit_comp_tot/2)
                    crit_cve_cost = crit_cve_eff * ap._day_rate /1000

                if high_comp_tot is None:
                    high_cve_eff = 0
                    high_cve_cost = 0
                else:
                    high_cve_eff = math.ceil(high_comp_tot/2)
                    high_cve_cost = high_cve_eff * ap._day_rate /1000

                if med_comp_tot is None:
                    med_cve_eff = 0
                    med_cve_cost = 0
                else:   
                    med_cve_eff = math.ceil(med_comp_tot/2)
                    med_cve_cost = med_cve_eff * ap._day_rate /1000

                summary_eff = int(summary_eff) + int(crit_cve_eff+high_cve_eff+med_cve_eff)
                summary_cost = round(summary_cost + crit_cve_cost + high_cve_cost + med_cve_cost,2)

                fix_now_eff = int(aip_fix_now_eff) + int(crit_cve_eff)
                fix_now_cost = round(aip_fix_now_cost + crit_cve_cost,2)

                near_term_eff = int(aip_near_term_eff) + int(high_cve_eff) + int(med_cve_eff)
                near_term_cost = round(aip_near_term_cost + high_cve_cost + med_cve_cost,2)

                self._ppt.replace_text(f'{{app{app_no}_high_sec_tot}}','{high_cve}')
                self._ppt.replace_text(f'{{app{app_no}_med_sec_tot}}','{mid__cve}')
                self._ppt.replace_text(f'{{app{app_no}_hl_fn_eff}}',crit_cve_eff)
                self._ppt.replace_text(f'{{app{app_no}_hl_nt_eff}}',int(high_cve_eff + med_cve_eff))


                '''
                    License compliance table
                
                '''
                lic_summary=self._hl_data.get_lic_info(hl_id)
                if not lic_summary.empty:
                    lic_summary=lic_summary[['component','version','release','license','risk']].drop_duplicates()
                    lic_summary.sort_values(['risk','license','component','version'],inplace=True)

                    lic_summary = lic_summary.groupby(['risk','license'])['component'].apply(lambda x: ','.join(x)).reset_index()
                    lic_summary['comp count']=lic_summary['component'].str.count(',')+1
                    
                    #remove duplicates but show full count
                    lic_summary['component']=lic_summary['component'].map(lambda x: no_dups(x,',',True))
                    
                    #only show the first 5 components
                    lic_summary['component']=lic_summary['component'].map(lambda x: x[:find_nth(x,',',6)])
                    lic_summary=lic_summary[['license','risk','comp count','component']]
                    lic_summary.sort_values(['risk','license','comp count'],inplace=True)

                    #remove low and undefined records from the table
                    lic_summary=lic_summary[lic_summary["risk"].str.contains("Low")==False]
                    lic_summary=lic_summary[lic_summary["risk"].str.contains("Undefined")==False]

                    #modify the forground color
                    lic_summary.loc[lic_summary['risk']=='High','forground']='255,0,0'
                    lic_summary.loc[lic_summary['risk']=='Medium','forground']='209,125,13'

                    #update the powerpoint table
                    self._ppt.update_table(f'app{app_no}_HL_table_lic_risks',
                                        lic_summary,include_index=False,
                                        forground='forground')
                
                    #add the high and medium license risk counts to the deck
                    self._ppt.replace_text(f'{{app{app_no}_high_lic_tot}}',
                        lic_summary[lic_summary['risk']=='High']['comp count'].sum())
                    self._ppt.replace_text(f'{{app{app_no}_med_lic_tot}}',
                        lic_summary[lic_summary['risk']=='Medium']['comp count'].sum())

                    self._ppt.update_table(f'app{app_no}_HL_table_CVEs',oss_df,include_index=False)



                """
                    Highlight risk slide
                """
                # cve_df = self._hl_data.get_cve_data(hl_id)[['component','cve','cweId','cweLabel','criticity']]
                # high_df = cve_df.groupby(['component'])


                # cve_df = cve_df.fillna('')
                # cve_df = cve_df.loc[cve_df['cweId'],'cweId']='NVD-CWE-noinfo'
                # cve_df.sort_values(by=['cweId','cve'], inplace=True, ascending = True)
                # risk_df = cve_df.groupby(['cweId','cve']).size

                """
                    Cloud ready excel sheet generation
                """
                try:
                    cloud = self._hl_data.get_cloud_info(hl_id)
                    cloud = cloud[['cloudRequirement.display','Technology','cloudRequirement.ruleType','cloudRequirement.criticality','contributionScore','roadblocks']]
                    file_name = f'{config.output}/cloud-{title_list[idx]}.xlsx'
                    writer = pd.ExcelWriter(file_name, engine='xlsxwriter')
                    col_widths=[50,10,10,10,10,10,10]
                    cloud_tab = util.format_table(writer,cloud,'Cloud Data',col_widths)
                    writer.save()
                except Exception as e:
                    self.error(f'unknown error while processing cloud ready data: {str(e)}')

            self._ppt.replace_text(f'{{app{app_no}_fn_tot_cost}}',fix_now_cost)
            self._ppt.replace_text(f'{{app{app_no}_fn_tot_eff}}',fix_now_eff)

            self._ppt.replace_text(f'{{app{app_no}_nt_tot_cost}}',near_term_cost)
            self._ppt.replace_text(f'{{app{app_no}_nt_tot_eff}}',near_term_eff)


            summary_total_cost = summary_total_cost + summary_cost
            self._ppt.replace_text(f'{{app{app_no}_summary_eff}}',summary_eff)
            self._ppt.replace_text(f'{{app{app_no}_summary_cost}}',summary_cost)
            self._ppt.replace_text(f'{{app{app_no}_summary_vio_cnt}}',summary_vio_cnt)
            self._ppt.replace_text(f'{{app{app_no}_summary_bus_txt}}',summary_bus_txt)
            self._ppt.replace_text(f'{{app{app_no}_summary_vio_txt}}',summary_vio_txt)

        self._ppt.replace_text('{summary_total_cost}',summary_total_cost)  
        if fix_now_eff > 0:
            show_stopper_flg = 'no'
        else:
            show_stopper_flg = ''
        self._ppt.replace_text('{show_stopper_flg}',show_stopper_flg)  

        avg_cost = float(summary_total_cost/app_cnt)   
        if avg_cost < 50:
            hml_cost_flg = 'low'
        elif avg_cost > 50 and avg_cost < 100:
            hml_cost_flg = 'medium'
        else:
            hml_cost_flg = 'high'
        self._ppt.replace_text('{hml_cost_flag}',hml_cost_flg)  

        self._ppt.replace_text('{daily_rate}',ap._day_rate)  

    def fill_critical_rules(self,app_id,app_no):
        rules_df = self._aip_data.critical_rules(app_id)
        if not rules_df.empty:
            critical_rule_df = pd.json_normalize(rules_df['rulePattern'])
            critical_rule_df = critical_rule_df[['name','critical']]
            rule_summary_df=critical_rule_df.groupby(['name']).size().reset_index(name='counts').sort_values(by=['counts'],ascending=False)
            rule_summary_df=rule_summary_df.head(5)
            self._ppt.update_table(f'app{app_no}_top_violations',rule_summary_df,include_index=False)

    def fill_violations(self,app_id,app_no):
        violation_df = pd.DataFrame(self._aip_data.violation_sizing(app_id),index=[0])
        violation_df['Violation Count']=pd.Series(["{0:,.0f}".format(val) for val in violation_df['Violation Count']])
        violation_df[' per file']=pd.Series(["{0:,.2f}".format(val) for val in violation_df[' per file']])
        violation_df[' per kLoC']=pd.Series(["{0:,.2f}".format(val) for val in violation_df[' per kLoC']])
        violation_df['Complex objects']=pd.Series(["{0:,.0f}".format(val) for val in violation_df['Complex objects']])
        violation_df[' With violations']=pd.Series(["{0:,.0f}".format(val) for val in violation_df[' With violations']])
        self._ppt.update_table(f'app{app_no}_violation_sizing',violation_df.transpose())
        self._ppt.replace_text(f'{{app{app_no}_critical_violations}}',violation_df['Violation Count'].loc[0])


    def fill_sizing(self,app_id,app_no):
        sizing_df = pd.DataFrame(self._aip_data.tech_sizing(app_id),index=[0])
        sizing_df['LoC']=pd.Series(["{0:,.0f} K".format(val / 1000) for val in sizing_df['LoC']])
        sizing_df['Files']=pd.Series(["{0:,.0f}".format(val) for val in sizing_df['Files']])
        sizing_df['Classes']=pd.Series(["{0:,.0f}".format(val) for val in sizing_df['Classes']])
        sizing_df['SQL Artifacts']=pd.Series(["{0:,.0f}".format(val) for val in sizing_df['SQL Artifacts']])
        sizing_df['Tables']=pd.Series(["{0:,.0f}".format(val) for val in sizing_df['Tables']])
        sizing_df = sizing_df.transpose()
        self._ppt.update_table(f'app{app_no}_tech_sizing',sizing_df)

if __name__ == '__main__':
    print('\nCAST Assessment Deck Generation Tool')
    print('Copyright (c) 2022 CAST Software Inc.\n')
    print('If you need assistance, please contact Nevin Kaplan (NKA) from the CAST US PS team\n')

    parser = argparse.ArgumentParser(description='Assessment Deck Generation Tool')
    parser.add_argument('-c','--config', required=True, help='Configuration properties file')
    args = parser.parse_args()
    ppt = GeneratePPT(Config(args.config))
    ppt.save_ppt()


